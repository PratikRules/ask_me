import streamlit as st
import os
from dotenv import load_dotenv
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import ebooklib.epub
from PIL import Image
import assemblyai as aai
from youtube_transcript_api import YouTubeTranscriptApi
from pytube import YouTube
import requests
from bs4 import BeautifulSoup
import pandas as pd
from pandasai import Agent
from pandasai.responses.response_parser import ResponseParser
import google.generativeai as genai
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_cohere import ChatCohere, CohereEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser

class PandasDataFrame(ResponseParser):
    def __init__(self, context) -> None:
        super().__init__(context)
    def format_dataframe(self, result):
        return result["value"]

load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
model1 = ChatCohere(model="command-r", temperature=0.9)
model2 = genai.GenerativeModel("gemini-1.0-pro-vision-latest")
model3 = ChatGoogleGenerativeAI(model="gemini-1.0-pro-latest", temperature=0.9)
embeddings = CohereEmbeddings(model="embed-english-v3.0")

def process_dirs():
    for dirs in ["dataframes", "images", "files", "faiss_index"]:
        if not os.path.exists(dirs):
            os.makedirs(dirs)
        else:
            for f in os.listdir(dirs):
                os.remove(os.path.join(dirs, f))

def document_path(doc):
    doc_path = os.path.join("files", doc.name)
    with open(doc_path, "wb") as f:
        f.write(doc.getvalue())
    return doc_path

def process_pdf(doc):
    text = ""
    pdf_reader = PdfReader(doc)
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n\n"
    st.session_state.file_type.append("Text")
    return text

def process_docx(doc):
    text = ""
    doc = Document(doc)
    for para in doc.paragraphs:
        text += para.text + "\n"
    st.session_state.file_type.append("Text")
    return text

def process_pptx(doc):
    text = ""
    prs = Presentation(doc)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text += run.text + "\n"
    st.session_state.file_type.append("Text")
    return text

def process_txt(doc):
    st.session_state.file_type.append("Text")
    return doc.read().decode("utf-8") + "\n"

def process_epub(doc):
    doc_path = document_path(doc)
    book = ebooklib.epub.read_epub(doc_path)
    text = ""
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.content, 'html.parser')
            text += soup.get_text() + "\n"
    st.session_state.file_type.append("Text")
    return text

def process_image(doc):
    image = Image.open(doc)
    image.save(f"images/{doc.name}")
    st.session_state.image.append(f"images/{doc.name}")
    st.session_state.file_type.append("Image")

def process_audio_video(doc, f_type):
    transcript = ""
    doc_path = document_path(doc)
    st.session_state.video.append(doc_path) if f_type == "Video" else st.session_state.audio.append(doc_path)
    transcriber = aai.Transcriber()
    transcript += transcriber.transcribe(doc_path).text
    st.session_state.file_type.append("Audio/Video")
    return transcript + "\n"

def process_csv_xlsx(doc, basename):
    df = pd.read_csv(doc) if doc.name.lower().endswith(".csv") else pd.read_excel(doc)
    st.session_state.spreadsheet.append(df)
    df.to_csv(f"dataframes/{basename}.csv", index=False)
    st.session_state.file_type.append("Spreadsheet")

def process_youtube_url(youtube_url):
    transcript_text = YouTubeTranscriptApi.get_transcript(video_id=youtube_url.split("=")[1])
    st.session_state.video.append(youtube_url)
    st.session_state.file_type.append("Audio/Video")
    transcript = ""
    for text_segment in transcript_text:
        transcript += " " + text_segment['text']
    transcript += f".\t{str(YouTube(youtube_url).vid_info['videoDetails'])}\n"
    return transcript

def process_web_url(web_url):
    with requests.get(web_url) as r:
        soup = BeautifulSoup(r.content, 'html.parser')
        st.session_state.file_type.append("Text")
        return f"\n\nWebsite Name: {soup.title.text}\n{soup.get_text()}"

def get_info(docs, any_url):
    text = ""
    process_dirs()
    for doc in docs:
        basename, extension = os.path.splitext(doc.name)
        if extension.lower() == ".pdf":
            text += f"\n\nPDF File Name: {basename}\n\n{process_pdf(doc)}"
        elif extension.lower() == ".docx":
            text += f"\n\nWord File Name: {basename}\n\n{process_docx(doc)}\n"
        elif extension.lower() == ".pptx":
            text += f"\n\nPowerPoint File Name: {basename}\n\n{process_pptx(doc)}\n"
        elif extension.lower() in [".txt", ".py"]:
            f_type = "Text" if extension.lower() == ".txt" else "Python"
            text += f"\n\n{f_type} File Name: {basename}\n\n{process_txt(doc)}\n"
        elif extension.lower() == ".epub":
            text += f"\n\nEpub File Name: {basename}\n\n{process_epub(doc)}\n"
        elif extension.lower() in [".jpg", ".jpeg", ".png"]:
            process_image(doc)
        elif extension.lower() in [".mp4", ".webm", ".mkv", ".mp3", ".wav"]:
            f_type = "Audio" if extension.lower() in [".mp3", ".wav"] else "Video"
            text += f"\n\n{f_type} File Name: {basename}\n\n{process_audio_video(doc, f_type)}\n"
        elif extension.lower() in [".csv", ".xlsx"]:
            process_csv_xlsx(doc, basename)
    if any_url:
        urls = any_url.split(",")
        for url in urls:
            if url.startswith("https://www.youtube.com/watch?v="):
                text += f"\n\nYoutube Video Name: {YouTube(url).title}\n\n{process_youtube_url(url)}"
            else:
                text += process_web_url(url)
    return text

def index(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    vector_store = FAISS.from_texts(chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def text_input(user_question):
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    retriever = new_db.as_retriever()
    prompt_template = """
    **As a proficient file reader, you excel in extracting information. Your task is to deliver a response closely
    aligned with both the context given in the user's question and the context within the file. Ensure accuracy and 
    detail in your answer, and keep the answer engaging. Highlight all the important words and keywords in bold.**

    **Context:**
    {context}

    **Chat History:**
    {chat_history}

    **Question:**
    {question}

    **Answer:**
    """
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "chat_history", "question"])
    chain = prompt | model1 | StrOutputParser()
    docs = retriever.invoke(user_question)
    response = chain.invoke({"context": docs, "question": user_question, "chat_history": st.session_state.memory})
    st.session_state.memory.extend([{"role": "user", "content": user_question}, {"role": "assistant", "content": response}])
    return response

def image_input(user_question):
    prompt = """
    . Write the answer in proper meaningful structured sentence(s). If possible, give more insights into the answer. 
    Write the answer in a proper format. If the answer is in a different language, translate it to English, give the 
    translated answer only, with proper meaningful structured sentence(s).
    """
    images = [user_question + prompt]
    for file in os.listdir("images"):
        images.append(Image.open(os.path.join("images", file)))
    response = model2.generate_content(images, generation_config={"temperature": 0.9})
    response.resolve()
    return response.text

def csv_input(user_question, agent):
    note = """
    . If the answer is a number, include them in a sentence as a string. If the question contains 'display',
    generate the appropriate dataframe. Only generate a plot, if the question mentions 'plot' or 'chart'.
    If a plot needs to be generated, include the following libraries only, in the following manner:
    ```python
    import pandas as pd
    import matplotlib.pyplot as plt
    import seaborn as sns
    ```
    It is important to use different color for each bar or category in the plot. Use hue and palette, in case of seaborn.
    For example: 'sns.barplot(x="variable", y="value", palette="Set2", data=dataframe)' Change the palette value at each instance.
    Rotate the x-axis labels by 45 degrees, if the x-axis labels are not numeric.
    Keep the figure size as 'plt.figure(figsize=(10, 6))'and also use 'plt.tight_layout()'.
    Make sure to use title, labels and other necessary parameters, to make the plot more presentable.
    Use 'plt.savefig('charts/name.png')' to save the plot, where name = Prompt ID
    """
    response = agent.chat(user_question + note)
    return response

def show_response(response):
    if isinstance(response, str) and ".png" in response:
        st.image(response)
    else:
        st.write(response)

def clear_display():
    for key in ['file_type', 'spreadsheet', 'image', 'audio', 'video']:
        setattr(st.session_state, key, [])
    st.session_state.processed = False
    st.session_state.agent = 1

def clear_chat_history():
    st.session_state.chat_history = [{"role": "assistant", "content": "How can I help you?"}]
    st.session_state.memory = []

def file_types():
    file_map = {
        "Text": (":orange[Text]", "Includes PDF, DOCX, PPTX, TXT"),
        "Spreadsheet": (":orange[Spreadsheet]", "Includes CSV, XLSX"),
        "Image": (":orange[Image]", "Includes PNG, JPG"),
        "Audio/Video": (":orange[Audio/Video]", "Includes MP3, WAV, MP4, WEBM, MKV"),
    }
    options, captions = [], []
    st.session_state.file_type = list(set(st.session_state.file_type))
    for file in st.session_state.file_type:
        if file in file_map:
            option, caption = file_map[file]
            options.append(option)
            captions.append(caption)
    return options, captions

def main():
    st.set_page_config("Ask Me", page_icon="📄", layout="wide")
    st.header("Ask Me", divider="orange")

    if "uploader_key" not in st.session_state:
        st.session_state.uploader_key = 0
    if all(key not in st.session_state for key in ["file_type", "processed", "spreadsheet", "agent", "image", "audio", "video"]):
        clear_display()
    if "chat_history" not in st.session_state:
        clear_chat_history()

    with st.sidebar:
        st.header("Menu", divider="orange")
        all_docs = st.file_uploader("Upload Files",
                                    type=["pdf", "docx", "pptx", "txt", "py", "csv", "xlsx", "png",
                                          "jpg", "jpeg", "mp3", "wav", "mp4", "webm", "mkv", "epub"],
                                    accept_multiple_files=True,
                                    key=f"file_uploader_{st.session_state.uploader_key}",
                                    label_visibility="collapsed")
        
        urls = None
        if st.toggle("Enter :orange[web] / :red[Youtube] URL(s)", key="url"):
            urls = st.text_area("Enter URL",
                                placeholder="Enter 1 or more URLs, separated by commas (Also works with YouTube)",
                                label_visibility="collapsed")

        if st.button("Submit and Process"):
            with st.spinner("Processing..."):
                clear_display()
                text = get_info(all_docs, urls)
                if text:
                    index(text)
                if st.session_state.file_type:
                    st.session_state.processed = True
                if "spreadsheet" in st.session_state:
                    csv_files = ["dataframes/" + f for f in os.listdir("dataframes")]
                    st.session_state.agent = Agent(csv_files, config={"llm": model3,
                                                    "save_charts": True,
                                                    "save_charts_path": "charts/",
                                                    "response_parser": PandasDataFrame})
            st.rerun()
        
        st.write("Select the type of file you want to ask")
        options, captions = file_types()
        file_type = st.radio("File Type",
                        options=options,
                        captions=captions,
                        disabled=not st.session_state.processed,
                        label_visibility="collapsed")
        
        col1, col2 = st.columns([3, 2])
        if col1.button("Clear chat history"):
            clear_chat_history()
        if col2.button(":red[Clear files]",
                       on_click=lambda: st.session_state.update({"url": False, "uploader_key": st.session_state.uploader_key + 1})):
            clear_display()
            process_dirs()
            st.rerun()
        
        for spreadsheets in st.session_state.spreadsheet:
            st.dataframe(spreadsheets, height=210, use_container_width=True)
        for images in st.session_state.image:
            st.image(images)
        for audios in st.session_state.audio:
            st.audio(audios)
        for videos in st.session_state.video:
            st.video(videos)
        
    for chat in st.session_state.chat_history:
        if chat["role"] == "user":
            with st.chat_message("user", avatar="user.png"):
                show_response(chat["content"])
        else:
            with st.chat_message("assistant", avatar="assistant.png"):
                show_response(chat["content"])

    user_question = st.chat_input("Ask a question from the files:", disabled=not st.session_state.processed)
    if user_question:
        st.chat_message("user", avatar="user.png").write(user_question)
        response = None
        with st.chat_message("assistant", avatar="assistant.png"):
            with st.spinner("Thinking..."):
                if file_type in [":orange[Text]", ":orange[Audio/Video]"]:
                    response = text_input(user_question)
                elif file_type == ":orange[Spreadsheet]":
                    response = csv_input(user_question, st.session_state.agent)
                else:
                    response = image_input(user_question)
                show_response(response)
        st.session_state.chat_history.extend([{"role": "user", "content": user_question}, {"role": "assistant", "content": response}])
        st.rerun()

if __name__ == "__main__":
    main()